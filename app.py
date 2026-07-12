# pip install streamlit langchain langchain-openai requests tavily-python python-docx python-pptx python-dotenv

import os
import io
import re
import json
import time
import requests
import streamlit as st
from dotenv import load_dotenv
from datetime import datetime, timedelta

# Langchain Imports
from langchain_openai import ChatOpenAI
from langchain.agents import tool, AgentExecutor, create_openai_tools_agent
from langchain_core.prompts import ChatPromptTemplate, MessagesPlaceholder
from langchain_community.tools.tavily_search import TavilySearchResults
# from langchain_core.tools import tool

# Document Export Imports
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Inches

# Load the .env file
load_dotenv()

# Set Tavily globally for LangChain
os.environ["TAVILY_API_KEY"] = os.getenv("TAVILY_API_KEY")

# ==========================================
# 1. INITIALIZE LLM & AGENT (Cached for performance)
# ==========================================
@st.cache_resource

def initialize_travel_agent():
    
    # Initialize LLM. Fetch from environment variables
    QWEN_API_KEY = os.getenv("QWEN_API_KEY")
    QWEN_BASE_URL = os.getenv("QWEN_BASE_URL")

    llm = ChatOpenAI(
    model="qwen-max", # Qwen-Max is best for complex reasoning and planning
    openai_api_key=QWEN_API_KEY,
    openai_api_base=QWEN_BASE_URL,
    temperature=0.7
)

    # Initialize LangChain's native Tavily tool
    # The LLM will automatically generate the search query based on the user's request
    tavily_search = TavilySearchResults(
        name="tavily_search_results",
        description="Useful for searching the web for real-time hotel recommendations, prices, and ratings.",
        max_results=5,
        search_depth="advanced",
        include_answer=True, # Gives the LLM a quick AI summary of the search
        include_raw_content=False
    )

    # --- External Tools ---
    @tool
    def get_current_date() -> str:
        """Get the current date to calculate trip start and end dates if the user doesn't specify exact dates."""
        return datetime.now().strftime("%Y-%m-%d")

    @tool
    def get_weather_forecast(city: str, start_date: str, end_date: str) -> str:
        """
        Get the daily weather forecast for a specific city and date range.
        Dates MUST be in YYYY-MM-DD format. Use this to avoid planning outdoor activities on rainy days.
        """
        try:
            geo_url = f"https://geocoding-api.open-meteo.com/v1/search?name={city}&count=1&language=en"
            geo_res = requests.get(geo_url).json()
            if not geo_res.get("results"):
                return f"Could not find coordinates for {city}."
            
            lat = geo_res["results"][0]["latitude"]
            lon = geo_res["results"][0]["longitude"]
            
            wx_url = (f"https://api.open-meteo.com/v1/forecast?latitude={lat}&longitude={lon}"
                    f"&daily=temperature_2m_max,temperature_2m_min,precipitation_probability_max"
                    f"&start_date={start_date}&end_date={end_date}&timezone=auto")
            wx_res = requests.get(wx_url).json()
            
            return json.dumps(wx_res.get("daily", {}))
        except Exception as e:
            return f"Weather API error: {str(e)}"

    @tool
    def get_travel_time_and_distance(origin_city: str, destination_city: str) -> str:
        """
        Calculate the realistic driving travel time (in minutes) and distance (in km) between two cities.
        Use this to ensure the itinerary is geographically logical and not too rushed.
        """
        try:
            headers = {'User-Agent': 'AITravelAgent/1.0 (Educational Project)'}
            
            def get_coords(city_name):
                url = f"https://nominatim.openstreetmap.org/search?q={city_name}&format=json&limit=1"
                res = requests.get(url, headers=headers).json()
                if res: return res[0]["lat"], res[0]["lon"]
                return None, None

            lat1, lon1 = get_coords(origin_city)
            time.sleep(1) # Respect Nominatim rate limits
            lat2, lon2 = get_coords(destination_city)
            
            if not all([lat1, lon1, lat2, lon2]):
                return "Could not find coordinates for one or both cities."

            route_url = f"http://router.project-osrm.org/route/v1/driving/{lon1},{lat1};{lon2},{lat2}?overview=false"
            route_res = requests.get(route_url).json()
            
            if route_res.get("code") == "Ok":
                route = route_res["routes"][0]
                distance_km = route["distance"] / 1000
                duration_mins = route["duration"] / 60
                return json.dumps({
                    "origin": origin_city, 
                    "destination": destination_city,
                    "distance_km": round(distance_km, 1), 
                    "duration_mins": round(duration_mins, 1)
                })
            return "Routing calculation failed."
        except Exception as e:
            return f"Maps API error: {str(e)}"
    
    # --- Sub-Agent Prompts & Executors ---
    # --- Itinerary Planner Agent ---
    planner_prompt = ChatPromptTemplate.from_messages([
        ("system", """You are an expert Itinerary Planner. Your job is to create a detailed day-by-day travel itinerary.
        
        STRICT RULES FOR PLANNING:
        1. Weather Adaptation: You MUST use the `get_weather_forecast` tool. If it rains heavily on a specific day, move outdoor/sightseeing activities to a sunny day and plan indoor activities (museums, cafes) for the rainy days.
        2. Geographical Logic: Use `get_travel_time_and_distance` to ensure daily activities are clustered together and not too rushed.
        3. Travel Style: Adapt activities strictly to the user's preferred style (e.g., relaxing, sightseeing, food).
        4. Output: Return ONLY the day-by-day itinerary in clean Markdown. Do NOT include hotel recommendations.
        CRITICAL LANGUAGE RULE: You MUST write your entire response strictly in English. Do not use Chinese or any other language under any circumstances.
        """),
        ("human", "{input}"),
        MessagesPlaceholder(variable_name="agent_scratchpad"),
    ])

    planner_agent = AgentExecutor(
        agent=create_openai_tools_agent(llm, [get_weather_forecast, get_travel_time_and_distance], planner_prompt),
        tools=[get_weather_forecast, get_travel_time_and_distance],
        verbose=True,
        handle_parsing_errors=True,
        max_iterations=5
    )

    # --- Hotel Recommender Agent ---
    hotel_prompt = ChatPromptTemplate.from_messages([
        ("system", """You are an expert Hotel Recommender. Your job is to find and recommend the best hotels.
        
        STRICT RULES:
        1. Use the `tavily_search_results` tool to find real hotels. Construct a highly specific search query including the city, star rating, and dates (e.g., "best 4-star hotels in Kyoto prices and ratings for Oct 10 to Oct 14").
        2. Read the search results and filter by the requested minimum star rating.
        3. Provide 3 top recommendations with names, estimated prices, ratings, and a brief reason why it fits the travel style.
        4. Output: Return ONLY the hotel recommendations in clean Markdown. Do NOT include itinerary details.
        CRITICAL LANGUAGE RULE: You MUST write your entire response strictly in English. Do not use Chinese or any other language under any circumstances.
        """),
        ("human", "{input}"),
        MessagesPlaceholder(variable_name="agent_scratchpad"),
    ])

    hotel_agent = AgentExecutor(
        agent=create_openai_tools_agent(llm, [tavily_search], hotel_prompt),
        tools=[tavily_search],
        verbose=True,
        handle_parsing_errors=True,
        max_iterations=3
    )

    # --- Food & Culinary Agent ---
    tavily_food_search = TavilySearchResults(
        name="tavily_food_search",
        description="Useful for searching the web for top-rated local food, must-try dishes, and highly-rated restaurants in a specific city.",
        max_results=5,
        search_depth="advanced",
        include_answer=True # Gives a quick AI summary of the best local dishes
    )

    food_prompt = ChatPromptTemplate.from_messages([
        ("system", """You are an expert Culinary Guide and Food Recommender. Your job is to find the best local food and highly-rated restaurants.
        
        STRICT RULES:
        1. Use the `tavily_food_search` tool to find must-try local dishes and top-rated restaurants. Construct highly specific queries like "must try local food and top rated restaurants in [City] with high google reviews and ratings".
        2. Focus strictly on places with high ratings (e.g., 4.5+ stars on Google Maps, TripAdvisor, or Yelp).
        3. Recommend 3-5 specific local dishes that are an absolute "must-try".
        4. Recommend 3 specific restaurants/cafes. For each, include the name, their signature dish, and mention their high ratings/reviews.
        5. Output: Return ONLY the food and restaurant recommendations in clean Markdown. Do NOT include itinerary or hotel details.
        CRITICAL LANGUAGE RULE: You MUST write your entire response strictly in English. Do not use Chinese or any other language under any circumstances.
        """),
        ("human", "{input}"),
        MessagesPlaceholder(variable_name="agent_scratchpad"),
    ])

    food_agent = AgentExecutor(
        agent=create_openai_tools_agent(llm, [tavily_food_search], food_prompt),
        tools=[tavily_food_search],
        verbose=True,
        handle_parsing_errors=True,
        max_iterations=3
    )

    # --- Supervisor Agent & Routing Tools ---
    @tool
    def plan_itinerary(task_description: str) -> str:
        """
        Delegate itinerary planning to the Itinerary Planner. 
        Pass a detailed description including city, start_date (YYYY-MM-DD), end_date (YYYY-MM-DD), length of stay, and travel style.
        """
        try:
            result = planner_agent.invoke({"input": task_description})
            return result["output"]
        except Exception as e:
            return f"Itinerary planning failed: {str(e)}"

    @tool
    def recommend_hotels(task_description: str) -> str:
        """
        Delegate hotel search to the Hotel Recommender. 
        Pass a detailed description including city, check_in (YYYY-MM-DD), check_out (YYYY-MM-DD), and minimum star rating.
        """
        try:
            result = hotel_agent.invoke({"input": task_description})
            return result["output"]
        except Exception as e:
            return f"Hotel recommendation failed: {str(e)}"

    @tool
    def recommend_food(task_description: str) -> str:
        """
        Delegate food and restaurant search to the Culinary Guide. 
        Pass a detailed description including the city and any specific dietary preferences or cuisine types requested.
        """
        try:
            result = food_agent.invoke({"input": task_description})
            return result["output"]
        except Exception as e:
            return f"Food recommendation failed: {str(e)}"

    # --- Supervisor Agent ---
    supervisor_prompt = ChatPromptTemplate.from_messages([
        ("system", """You are the Elite Travel Concierge Supervisor. 
        Your job is to take the user's travel request, delegate tasks to the Itinerary Planner, Hotel Recommender, and Culinary Guide, and synthesize their outputs into a final, beautiful response.
        
        WORKFLOW:
        1. Date Calculation: If the user didn't specify exact dates, use `get_current_date` and assume the trip starts tomorrow. Calculate the end date based on the length of stay.
        2. Delegation 1: Call the `plan_itinerary` tool with the city, start_date, end_date, length of stay, and travel style.
        3. Delegation 2: Call the `recommend_hotels` tool with the city, check_in, check_out, and minimum star rating.
        4. Delegation 3: Call the `recommend_food` tool with the city and any food preferences.
        
        SYNTHESIS & FORMATTING RULES:
        - Combine all three delegated results into a single, cohesive response.
        - Use clear Markdown headers: `## 🗓️ Day-by-Day Itinerary`, `## 🏨 Hotel Recommendations`, and `## 🍜 Must-Try Food & Dining`.
        - Use bold text, bullet points, and emojis to make it visually appealing and easy to read.
        - Add a polite, enthusiastic intro and a helpful outro (e.g., reminding them to book in advance).
        - Do NOT invent information. Only use the data returned by your sub-agents.
        
        CRITICAL LANGUAGE RULE: Your final output MUST be 100% in English. If any sub-agent returns text in Chinese or another language, you MUST translate it into English before presenting the final response to the user.
        """),
        ("human", "{input}"),
        MessagesPlaceholder(variable_name="agent_scratchpad"),
    ])

    # Note: We added 'recommend_food' to the tools list and increased max_iterations to 8 
    # because the supervisor now has to make 4 tool calls (date, itinerary, hotel, food) before synthesizing.
    return AgentExecutor(
        agent=create_openai_tools_agent(
            llm, 
            [get_current_date, plan_itinerary, recommend_hotels, recommend_food], 
            supervisor_prompt
        ),
        tools=[get_current_date, plan_itinerary, recommend_hotels, recommend_food],
        verbose=True,
        handle_parsing_errors=True,
        max_iterations=8 
    )

# ==========================================
# 2. DOCUMENT GENERATION HELPERS
# ==========================================
def generate_docx(markdown_text):
    doc = Document()
    doc.add_heading('Travel Itinerary & Recommendations', 0)
    
    for line in markdown_text.split('\n'):
        line = line.strip()
        if not line: continue
        
        if line.startswith('## '):
            doc.add_heading(line.replace('## ', ''), level=1)
        elif line.startswith('### '):
            doc.add_heading(line.replace('### ', ''), level=2)
        elif line.startswith('- ') or line.startswith('* '):
            doc.add_paragraph(line[2:], style='List Bullet')
        else:
            # Basic bold parsing
            clean_line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
            doc.add_paragraph(clean_line)
            
    byte_io = io.BytesIO()
    doc.save(byte_io)
    byte_io.seek(0)
    return byte_io

def generate_pptx(markdown_text):
    prs = Presentation()
    
    # Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Your Elite Travel Plan"
    title_slide.placeholders[1].text = "Generated by AI Travel Concierge"

    # Split by H2 headers for new slides
    sections = re.split(r'(?=^## )', markdown_text, flags=re.MULTILINE)
    
    for section in sections:
        if not section.strip(): continue
        
        lines = section.strip().split('\n')
        slide_title = lines[0].replace('## ', '')
        slide_body = '\n'.join(lines[1:]).strip()
        
        # Content Slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title
        
        # Clean up text for PPTX (remove markdown bold/italics)
        clean_body = re.sub(r'\*\*(.*?)\*\*', r'\1', slide_body)
        clean_body = re.sub(r'\*(.*?)\*', r'\1', clean_body)
        clean_body = clean_body.replace('- ', '• ')
        
        slide.placeholders[1].text = clean_body

    byte_io = io.BytesIO()
    prs.save(byte_io)
    byte_io.seek(0)
    return byte_io

# ==========================================
# 3. STREAMLIT FRONT-END UI
# ==========================================
st.set_page_config(page_title="Elite Travel Concierge", page_icon="✈️", layout="wide")

st.title("✈️ Elite Travel Concierge")
st.markdown("Plan your perfect trip with AI. Select your preferences below and let our agents craft a personalised itinerary, hotel and food guide.")

# --- Sidebar / Input Form ---
with st.form("travel_form"):
    st.subheader("🌍 Trip Details")
    
    col1, col2 = st.columns(2)
    with col1:
        destination = st.text_input("Destination Country / City", placeholder="e.g., South Korea (Seoul & Jeju)")
        hotel_stars = st.select_slider("Minimum Hotel Star Rating", options=[1, 2, 3, 4, 5], value=1)
        
    with col2:
        travel_dates = st.date_input(
            "Travel Dates (Select Start and End Date)", 
            value=[], 
            min_value=datetime.now().date(),
            format="YYYY-MM-DD"
        )
        travel_style = st.multiselect(
            "Travel Style / Preferences", 
            ["Sightseeing", "Foodie", "Shopping", "Relaxing", "Adventure", "Culture"],
            default=["Sightseeing", "Foodie"]
        )

    submitted = st.form_submit_button("✨ Generate Travel Plan", use_container_width=True)

# --- Execution Logic ---
if submitted:
    if not destination:
        st.error("Please enter a destination.")
    elif len(travel_dates) != 2:
        st.error("Please select both a start and end date.")
    else:
        start_date = travel_dates[0].strftime("%Y-%m-%d")
        end_date = travel_dates[1].strftime("%Y-%m-%d")
        days = (travel_dates[1] - travel_dates[0]).days + 1
        styles = ", ".join(travel_style)
        
        user_prompt = (f"I am going to {destination} for {days} days from {start_date} to {end_date}. "
                       f"I prefer a {styles} style. I need at least {hotel_stars}-star hotels. "
                       f"Please provide the entire travel plan in English.")
        
        st.info(f"**Prompt sent to Agent:** {user_prompt}")
        
        with st.spinner("🤖 Our AI Agents are planning your trip and customising your travel plan ... (This may take 30-60 seconds)"):
            agent = initialize_travel_agent()
            result = agent.invoke({"input": user_prompt})
            final_output = result["output"]
            
            # Save to session state so it persists on UI rerenders
            st.session_state['travel_plan'] = final_output

# --- Display Output and Downloads ---
if 'travel_plan' in st.session_state:
    st.divider()
    st.subheader("📝 Your Personalised Travel Plan")
    
    # Display Markdown
    st.markdown(st.session_state['travel_plan'])
    
    st.divider()
    st.subheader("📥 Download Your Travel Plan")
    
    col1, col2 = st.columns(2)
    
    with col1:
        docx_file = generate_docx(st.session_state['travel_plan'])
        st.download_button(
            label="📄 Download as Word (.docx)",
            data=docx_file,
            file_name="Travel_Itinerary.docx",
            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            use_container_width=True
        )
        
    with col2:
        pptx_file = generate_pptx(st.session_state['travel_plan'])
        st.download_button(
            label="📊 Download as PowerPoint (.pptx)",
            data=pptx_file,
            file_name="Travel_Itinerary.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            use_container_width=True
        )